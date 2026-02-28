"""
DocAI Platform - 文档解析模块
支持 PDF / DOCX / DOC / PPTX / XLSX / CSV / TXT / MD
解析策略路由 → 提取文本 + 保留结构（标题层级、页码、表格）
"""

from __future__ import annotations

import csv
import io
import os
import re
from pathlib import Path

import structlog

from app.core.models import ParsedDocument, Section, TableData

logger = structlog.get_logger()


class UnsupportedFormatError(Exception):
    """不支持的文档格式"""
    pass


# ═══════════════════════════════════════════════════════════════════════════
# 主入口
# ═══════════════════════════════════════════════════════════════════════════

def parse_document(file_path: str) -> ParsedDocument:
    """
    根据文件扩展名选择合适的解析器

    Args:
        file_path: 本地文件路径

    Returns:
        ParsedDocument: 结构化的文档解析结果
    """
    ext = Path(file_path).suffix.lower()
    filename = Path(file_path).name

    logger.info("Parsing document", file=filename, ext=ext)

    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".doc": _parse_doc,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".csv": _parse_csv,
        ".txt": _parse_text,
        ".md": _parse_text,
    }

    parser = parsers.get(ext)
    if parser is None:
        raise UnsupportedFormatError(f"Unsupported file format: {ext}")

    doc = parser(file_path)
    doc.filename = filename

    # 如果没有 title，用文件名
    if not doc.title:
        doc.title = Path(filename).stem

    logger.info(
        "Document parsed",
        file=filename,
        sections=len(doc.sections),
        tables=len(doc.tables),
        pages=doc.page_count,
    )
    return doc


# ═══════════════════════════════════════════════════════════════════════════
# PDF 解析
# ═══════════════════════════════════════════════════════════════════════════

def _is_scanned_pdf(file_path: str) -> bool:
    """检测 PDF 是否为扫描件（文字很少则认为是扫描件）"""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    total_text = 0
    sample_pages = min(3, len(doc))  # 仅检查前 3 页
    for i in range(sample_pages):
        page = doc[i]
        text = page.get_text()
        total_text += len(text.strip())
    doc.close()

    # 平均每页少于 50 个字符则认为是扫描件
    avg_chars = total_text / max(sample_pages, 1)
    return avg_chars < 50


def _parse_pdf(file_path: str) -> ParsedDocument:
    """使用 PyMuPDF 解析 PDF"""
    import fitz

    if _is_scanned_pdf(file_path):
        logger.warning("Scanned PDF detected, OCR not available in Phase 1", file=file_path)
        # Phase 1 先不支持 OCR，返回空白
        return ParsedDocument(
            title=Path(file_path).stem,
            page_count=0,
            raw_text="[Scanned PDF - OCR processing required]",
        )

    doc = fitz.open(file_path)
    page_count = len(doc)

    # 提取目录结构（如果有）
    toc = doc.get_toc()  # [[level, title, page], ...]

    all_sections: list[Section] = []
    all_tables: list[TableData] = []
    raw_text_parts: list[str] = []

    if toc:
        # 有目录结构的 PDF：按目录切分章节
        all_sections = _build_sections_from_toc(doc, toc)
    else:
        # 无目录：按页面切分，尝试识别标题
        all_sections = _build_sections_from_pages(doc)

    # 提取表格
    for page_idx in range(page_count):
        page = doc[page_idx]
        raw_text_parts.append(page.get_text())

        # PyMuPDF 表格提取
        try:
            tables = page.find_tables()
            for table in tables:
                md_table = _fitz_table_to_markdown(table)
                if md_table:
                    all_tables.append(TableData(
                        content=md_table,
                        page_number=page_idx + 1,
                    ))
        except Exception:
            pass  # 部分页面可能没有表格

    doc.close()

    return ParsedDocument(
        title=Path(file_path).stem,
        page_count=page_count,
        sections=all_sections,
        tables=all_tables,
        raw_text="\n".join(raw_text_parts),
    )


def _build_sections_from_toc(doc, toc: list) -> list[Section]:
    """根据 PDF 目录 (TOC) 构建章节结构"""
    import fitz

    sections = []
    page_count = len(doc)

    for i, (level, title, page_num) in enumerate(toc):
        # 确定章节结束页（下一个同级或更高级章节的开始页）
        end_page = page_count
        for j in range(i + 1, len(toc)):
            if toc[j][0] <= level:
                end_page = toc[j][2]
                break

        # 提取章节文本
        content_parts = []
        page_numbers = []
        for p in range(max(0, page_num - 1), min(end_page, page_count)):
            page = doc[p]
            text = page.get_text().strip()
            if text:
                content_parts.append(text)
                page_numbers.append(p + 1)

        content = "\n".join(content_parts)
        # 去掉标题本身（如果内容以标题开头）
        if content.startswith(title):
            content = content[len(title):].strip()

        sections.append(Section(
            title=title,
            level=level,
            content=content,
            page_numbers=sorted(set(page_numbers)),
        ))

    return sections


def _build_sections_from_pages(doc) -> list[Section]:
    """没有 TOC 的 PDF，按页面分段，尝试识别标题"""
    sections = []
    current_text = []
    current_pages = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        text = page.get_text().strip()
        if not text:
            continue

        current_text.append(text)
        current_pages.append(page_idx + 1)

        # 每 3 页或最后一页生成一个 section
        if len(current_text) >= 3 or page_idx == len(doc) - 1:
            full_text = "\n".join(current_text)
            # 尝试用第一行作为标题
            lines = full_text.split("\n")
            title = lines[0][:100] if lines else f"Page {current_pages[0]}"
            content = "\n".join(lines[1:]) if len(lines) > 1 else ""

            sections.append(Section(
                title=title,
                level=1,
                content=content,
                page_numbers=list(current_pages),
            ))
            current_text = []
            current_pages = []

    return sections


def _fitz_table_to_markdown(table) -> str:
    """将 PyMuPDF 表格对象转为 Markdown"""
    try:
        data = table.extract()
        if not data or len(data) < 2:
            return ""

        # 第一行作为表头
        header = data[0]
        header = [str(cell).strip() if cell else "" for cell in header]
        separator = ["-" * max(3, len(h)) for h in header]
        rows = []
        for row in data[1:]:
            cells = [str(cell).strip() if cell else "" for cell in row]
            rows.append(cells)

        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in rows:
            # 确保列数匹配
            while len(row) < len(header):
                row.append("")
            lines.append("| " + " | ".join(row[:len(header)]) + " |")

        return "\n".join(lines)
    except Exception:
        return ""


# ═══════════════════════════════════════════════════════════════════════════
# DOCX 解析
# ═══════════════════════════════════════════════════════════════════════════

def _parse_docx(file_path: str) -> ParsedDocument:
    """使用 python-docx 解析 DOCX"""
    from docx import Document as DocxDocument
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

    doc = DocxDocument(file_path)

    sections: list[Section] = []
    tables: list[TableData] = []
    raw_parts: list[str] = []

    # 追踪当前章节栈
    section_stack: list[Section] = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # 段落
            para = None
            for p in doc.paragraphs:
                if p._element is element:
                    para = p
                    break
            if para is None:
                continue

            text = para.text.strip()
            if not text:
                continue

            raw_parts.append(text)

            # 检测是否为标题
            style_name = para.style.name if para.style else ""
            heading_level = _get_heading_level(style_name, text)

            if heading_level > 0:
                new_section = Section(
                    title=text,
                    level=heading_level,
                    content="",
                )
                # 找到合适的父级
                while section_stack and section_stack[-1].level >= heading_level:
                    section_stack.pop()

                if section_stack:
                    section_stack[-1].children.append(new_section)
                else:
                    sections.append(new_section)
                section_stack.append(new_section)
            else:
                # 普通段落，追加到当前章节
                if section_stack:
                    if section_stack[-1].content:
                        section_stack[-1].content += "\n" + text
                    else:
                        section_stack[-1].content = text
                else:
                    # 没有章节的文档，创建一个默认章节
                    default_section = Section(
                        title="",
                        level=0,
                        content=text,
                    )
                    sections.append(default_section)
                    section_stack.append(default_section)

        elif tag == "tbl":
            # 表格
            for tbl in doc.tables:
                if tbl._element is element:
                    md_table = _docx_table_to_markdown(tbl)
                    if md_table:
                        section_path = ""
                        if section_stack:
                            section_path = section_stack[-1].title
                        tables.append(TableData(
                            content=md_table,
                            section_path=section_path,
                        ))
                        raw_parts.append(md_table)
                    break

    # 如果没有识别到章节结构，按固定段落数分块
    if not sections or (len(sections) == 1 and not sections[0].title):
        sections = _fallback_sections_from_text("\n".join(raw_parts))

    return ParsedDocument(
        title=Path(file_path).stem,
        sections=_flatten_sections(sections),
        tables=tables,
        raw_text="\n".join(raw_parts),
    )


def _get_heading_level(style_name: str, text: str) -> int:
    """检测段落是否为标题，返回标题层级 (0=非标题)"""
    # 方式 1：Word 内置标题样式
    if style_name.startswith("Heading"):
        try:
            return int(style_name.replace("Heading", "").strip())
        except ValueError:
            return 1

    # 中文标题样式
    cn_heading_map = {
        "标题 1": 1, "标题 2": 2, "标题 3": 3, "标题 4": 4,
        "标题1": 1, "标题2": 2, "标题3": 3, "标题4": 4,
    }
    if style_name in cn_heading_map:
        return cn_heading_map[style_name]

    # 方式 2：通过文本模式检测（如 "1.2.3 ..."、"第一章"、"一、"）
    heading_patterns = [
        (r"^第[一二三四五六七八九十百]+[章节部分篇]", 1),
        (r"^[一二三四五六七八九十]+[、．.]", 1),
        (r"^(\d+)\s+\S", 1),
        (r"^(\d+\.\d+)\s+\S", 2),
        (r"^(\d+\.\d+\.\d+)\s+\S", 3),
        (r"^\(([一二三四五六七八九十]+)\)", 2),
        (r"^（([一二三四五六七八九十]+)）", 2),
    ]

    for pattern, level in heading_patterns:
        if re.match(pattern, text) and len(text) < 100:
            return level

    return 0


def _docx_table_to_markdown(table) -> str:
    """将 python-docx 表格转为 Markdown"""
    try:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        header = rows[0]
        separator = ["-" * max(3, len(h)) for h in header]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in rows[1:]:
            while len(row) < len(header):
                row.append("")
            lines.append("| " + " | ".join(row[:len(header)]) + " |")

        return "\n".join(lines)
    except Exception:
        return ""


def _flatten_sections(sections: list[Section]) -> list[Section]:
    """将嵌套的章节结构展平为列表（保留 section_path 信息）"""
    result = []

    def _flatten(sec: Section, parent_path: str = ""):
        path = sec.get_section_path(parent_path)
        # 创建一个新的 section，设置 section_path 信息
        flat = Section(
            title=sec.title,
            level=sec.level,
            content=sec.content,
            page_numbers=sec.page_numbers,
        )
        result.append(flat)

        for child in sec.children:
            _flatten(child, path)

    for s in sections:
        _flatten(s)

    return result


def _fallback_sections_from_text(text: str, lines_per_section: int = 30) -> list[Section]:
    """fallback: 无法识别章节时，按固定行数分段"""
    lines = text.split("\n")
    sections = []
    for i in range(0, len(lines), lines_per_section):
        group = lines[i:i + lines_per_section]
        content = "\n".join(group)
        title_line = group[0][:80] if group else f"Segment {i // lines_per_section + 1}"
        sections.append(Section(
            title=title_line,
            level=1,
            content=content,
        ))
    return sections


# ═══════════════════════════════════════════════════════════════════════════
# DOC 解析 (旧版 Word，使用 docling 或 textract fallback)
# ═══════════════════════════════════════════════════════════════════════════

def _parse_doc(file_path: str) -> ParsedDocument:
    """解析旧版 .doc 文件"""
    # 尝试使用 docling
    try:
        return _parse_with_docling(file_path)
    except Exception as e:
        logger.warning("Docling failed for .doc, trying antiword", error=str(e))

    # fallback: 用 antiword 或其他系统命令
    try:
        import subprocess
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", file_path],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0 and result.stdout.strip():
            text = result.stdout
            sections = _fallback_sections_from_text(text)
            return ParsedDocument(
                title=Path(file_path).stem,
                sections=sections,
                raw_text=text,
            )
    except Exception:
        pass

    # 最后 fallback: 用 docling
    return _parse_with_docling(file_path)


# ═══════════════════════════════════════════════════════════════════════════
# PPTX 解析
# ═══════════════════════════════════════════════════════════════════════════

def _parse_pptx(file_path: str) -> ParsedDocument:
    """解析 PPT 文件"""
    from pptx import Presentation

    prs = Presentation(file_path)
    sections: list[Section] = []
    tables: list[TableData] = []
    page_count = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        slide_title = ""

        for shape in slide.shapes:
            # 提取标题
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if shape.shape_type is not None and "TITLE" in str(shape.shape_type):
                        slide_title = text
                    else:
                        slide_texts.append(text)

            # 提取表格
            if shape.has_table:
                md_table = _pptx_table_to_markdown(shape.table)
                if md_table:
                    tables.append(TableData(
                        content=md_table,
                        page_number=slide_idx,
                        section_path=f"Slide {slide_idx}",
                    ))
                    slide_texts.append(md_table)

        if not slide_title:
            slide_title = f"Slide {slide_idx}"

        sections.append(Section(
            title=slide_title,
            level=1,
            content="\n".join(slide_texts),
            page_numbers=[slide_idx],
        ))

    return ParsedDocument(
        title=Path(file_path).stem,
        page_count=page_count,
        sections=sections,
        tables=tables,
        raw_text="\n\n".join(
            f"## {s.title}\n{s.content}" for s in sections
        ),
    )


def _pptx_table_to_markdown(table) -> str:
    """将 PPTX 表格转为 Markdown"""
    try:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        header = rows[0]
        separator = ["-" * max(3, len(h)) for h in header]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in rows[1:]:
            while len(row) < len(header):
                row.append("")
            lines.append("| " + " | ".join(row[:len(header)]) + " |")

        return "\n".join(lines)
    except Exception:
        return ""


# ═══════════════════════════════════════════════════════════════════════════
# Excel / CSV 解析
# ═══════════════════════════════════════════════════════════════════════════

def _parse_xlsx(file_path: str) -> ParsedDocument:
    """解析 Excel 文件——每个 Sheet 作为一个 section，数据转 Markdown 表格"""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    sections: list[Section] = []
    tables: list[TableData] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            # 跳过完全空行
            if any(c for c in cells):
                rows.append(cells)

        if not rows:
            sections.append(Section(
                title=sheet_name,
                level=1,
                content="(empty sheet)",
            ))
            continue

        # 转 Markdown 表格
        header = rows[0]
        separator = ["-" * max(3, len(str(h))) for h in header]
        md_lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in rows[1:]:
            while len(row) < len(header):
                row.append("")
            md_lines.append("| " + " | ".join(row[:len(header)]) + " |")

        md_table = "\n".join(md_lines)

        tables.append(TableData(
            content=md_table,
            section_path=sheet_name,
            caption=f"Sheet: {sheet_name}",
        ))

        sections.append(Section(
            title=sheet_name,
            level=1,
            content=md_table,
        ))

    wb.close()

    return ParsedDocument(
        title=Path(file_path).stem,
        page_count=len(wb.sheetnames),
        sections=sections,
        tables=tables,
        raw_text="\n\n".join(s.content for s in sections),
    )


def _parse_csv(file_path: str) -> ParsedDocument:
    """解析 CSV 文件"""
    with open(file_path, "r", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return ParsedDocument(title=Path(file_path).stem)

    header = rows[0]
    separator = ["-" * max(3, len(h)) for h in header]
    md_lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in rows[1:]:
        while len(row) < len(header):
            row.append("")
        md_lines.append("| " + " | ".join(row[:len(header)]) + " |")

    md_table = "\n".join(md_lines)

    return ParsedDocument(
        title=Path(file_path).stem,
        sections=[Section(title="Data", level=1, content=md_table)],
        tables=[TableData(content=md_table, caption=Path(file_path).name)],
        raw_text=md_table,
    )


# ═══════════════════════════════════════════════════════════════════════════
# Text / Markdown 解析
# ═══════════════════════════════════════════════════════════════════════════

def _parse_text(file_path: str) -> ParsedDocument:
    """解析纯文本或 Markdown 文件"""
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    ext = Path(file_path).suffix.lower()

    if ext == ".md":
        sections = _parse_markdown_sections(text)
    else:
        sections = _fallback_sections_from_text(text)

    return ParsedDocument(
        title=Path(file_path).stem,
        sections=sections,
        raw_text=text,
    )


def _parse_markdown_sections(text: str) -> list[Section]:
    """解析 Markdown 的标题结构"""
    sections = []
    current_section: Section | None = None
    content_lines = []

    for line in text.split("\n"):
        heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
        if heading_match:
            # 保存之前的 section
            if current_section is not None:
                current_section.content = "\n".join(content_lines).strip()
                sections.append(current_section)
                content_lines = []

            level = len(heading_match.group(1))
            title = heading_match.group(2).strip()
            current_section = Section(title=title, level=level)
        else:
            content_lines.append(line)

    # 保存最后一个 section
    if current_section is not None:
        current_section.content = "\n".join(content_lines).strip()
        sections.append(current_section)
    elif content_lines:
        sections.append(Section(
            title="",
            level=0,
            content="\n".join(content_lines).strip(),
        ))

    return sections


# ═══════════════════════════════════════════════════════════════════════════
# Docling 解析 (通用 fallback)
# ═══════════════════════════════════════════════════════════════════════════

def _parse_with_docling(file_path: str) -> ParsedDocument:
    """使用 docling 解析文档（支持多种格式）"""
    from docling.document_converter import DocumentConverter

    converter = DocumentConverter()
    result = converter.convert(file_path)

    # docling 返回的主要内容
    md_text = result.document.export_to_markdown()

    sections = _parse_markdown_sections(md_text) if md_text else []

    return ParsedDocument(
        title=Path(file_path).stem,
        sections=sections,
        raw_text=md_text,
    )
